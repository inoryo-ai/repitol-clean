# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

DN=RGBColor(15,23,42); IND=RGBColor(79,70,229); SL=RGBColor(100,116,139)
LBG=RGBColor(241,245,249); WH=RGBColor(255,255,255); GR=RGBColor(5,150,105)
WN=RGBColor(217,119,6); ER=RGBColor(220,38,38); CVR=RGBColor(30,41,59)
BD=RGBColor(226,232,240); SIND=RGBColor(129,140,248); HBG=RGBColor(237,233,254)
FT='Meiryo'; M=Inches(1.2); CW=Inches(10.933)

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5); SW=prs.slide_width

def bg(s,c): s.background.fill.solid(); s.background.fill.fore_color.rgb=c
def rc(s,l,t,w,h,f=None,ln=None,lw=Pt(1)):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,l,t,w,h)
    if f: sh.fill.solid(); sh.fill.fore_color.rgb=f
    else: sh.fill.background()
    if ln: sh.line.color.rgb=ln; sh.line.width=lw
    else: sh.line.fill.background()
    return sh
def tx(s,l,t,w,h,txt,sz=14,c=DN,b=False,a=PP_ALIGN.LEFT):
    tb=s.shapes.add_textbox(l,t,w,h); tf=tb.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.text=txt; p.font.size=Pt(sz); p.font.color.rgb=c; p.font.bold=b; p.font.name=FT; p.alignment=a
    return tb
def ap(tf,t,sz=14,c=DN,b=False,a=PP_ALIGN.LEFT):
    p=tf.add_paragraph(); p.text=t; p.font.size=Pt(sz); p.font.color.rgb=c; p.font.bold=b; p.font.name=FT; p.alignment=a
def hd(s,n,t):
    bg(s,WH); rc(s,M,Inches(0.6),CW,Pt(3),f=IND)
    bd=rc(s,M,Inches(0.75),Inches(0.5),Inches(0.5),f=IND)
    p=bd.text_frame.paragraphs[0]; p.text=n; p.font.size=Pt(16); p.font.color.rgb=WH; p.font.bold=True; p.font.name=FT; p.alignment=PP_ALIGN.CENTER
    tx(s,Inches(1.9),Inches(0.72),Inches(9.5),Inches(0.55),t,sz=26,c=DN,b=True)

# 1. COVER
s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s,CVR)
rc(s,Inches(5.667),Inches(0.8),Inches(2.0),Inches(2.0),f=WH)
tx(s,Inches(5.667),Inches(1.0),Inches(2.0),Inches(1.6),'\U0001f3af',sz=55,c=DN,a=PP_ALIGN.CENTER)
tx(s,Inches(0),Inches(3.1),SW,Inches(0.7),'SlideShot',sz=48,c=WH,b=True,a=PP_ALIGN.CENTER)
tx(s,Inches(0),Inches(3.8),SW,Inches(0.4),'AI\u55b6\u696d\u30c9\u30ad\u30e5\u30e1\u30f3\u30c8\u81ea\u52d5\u751f\u6210 SaaS',sz=18,c=SIND,a=PP_ALIGN.CENTER)
rc(s,Inches(4.0),Inches(4.5),Inches(5.333),Pt(1),f=RGBColor(80,80,120))
tx(s,Inches(2.5),Inches(4.8),Inches(8.333),Inches(0.6),'\u5165\u529b\u3059\u308b\u3060\u3051\u3067\u3001\u55b6\u696d\u30e1\u30fc\u30eb\u30fb\u63d0\u6848\u66f8\u30fb\u898b\u7a4d\u66f8\u304c\u5b8c\u6210',sz=20,c=RGBColor(210,210,230),a=PP_ALIGN.CENTER)
tx(s,Inches(0),Inches(6.5),SW,Inches(0.35),'\u682a\u5f0f\u4f1a\u793e Centaurus  |  2026\u5e743\u6708',sz=13,c=RGBColor(140,140,170),a=PP_ALIGN.CENTER)

# 2. PROBLEM
s=prs.slides.add_slide(prs.slide_layouts[6]); hd(s,'01','\u8ab2\u984c \u2015 \u55b6\u696d\u73fe\u5834\u306e\u30da\u30a4\u30f3')
probs=[('\u30e1\u30fc\u30eb\u4f5c\u6210\u306b\u6642\u9593\u304c\u304b\u304b\u308b','1\u901a\u3042\u305f\u308a\u5e73\u571215\u301c30\u5206\u3002\n\u696d\u7a2e\u30fb\u76f8\u624b\u306b\u5fdc\u3058\u305f\u8abf\u6574\u304c\u5fc5\u8981\u3002'),
('\u63d0\u6848\u66f8\u304c\u5c5e\u4eba\u5316',':\u500b\u4eba\u306e\u30b9\u30ad\u30eb\u306b\u4f9d\u5b58\u3057\u3001\n\u54c1\u8cea\u306b\u30d0\u30e9\u3064\u304d\u304c\u3042\u308b\u3002'),
('\u898b\u7a4d\u66f8\u304c\u30d0\u30e9\u30d0\u30e9','Excel\u624b\u4f5c\u696d\u3067\u4f5c\u6210\u3002\n\u8ee2\u8a18\u30df\u30b9\u3084\u629c\u3051\u6f0f\u308c\u304c\u767a\u751f\u3002'),
('\u30c6\u30f3\u30d7\u30ec\u304c\u6700\u9069\u5316\u3055\u308c\u306a\u3044','\u3069\u306e\u6587\u9762\u304c\u53d7\u6ce8\u306b\u7e4b\u304c\u3063\u305f\u304b\n\u30c7\u30fc\u30bf\u304c\u84c4\u7a4d\u3055\u308c\u306a\u3044\u3002')]
for i,(t,d) in enumerate(probs):
    x=M+Inches((i%2)*5.6); y=Inches(1.5)+Inches((i//2)*2.2)
    rc(s,x,y,Inches(5.0),Inches(1.8),f=LBG,ln=BD); rc(s,x,y,Pt(5),Inches(1.8),f=ER)
    tx(s,x+Inches(0.25),y+Inches(0.15),Inches(4.5),Inches(0.35),'\u2716 '+t,sz=16,c=ER,b=True)
    tx(s,x+Inches(0.25),y+Inches(0.6),Inches(4.5),Inches(1.0),d,sz=14,c=SL)

# 3. SOLUTION
s=prs.slides.add_slide(prs.slide_layouts[6]); hd(s,'02','\u30bd\u30ea\u30e5\u30fc\u30b7\u30e7\u30f3')
rc(s,M,Inches(1.5),CW,Inches(1.0),f=HBG,ln=IND,lw=Pt(2))
tx(s,M+Inches(0.3),Inches(1.6),CW-Inches(0.6),Inches(0.8),'\u55b6\u696d\u30e1\u30fc\u30eb\u30fb\u63d0\u6848\u66f8\u30fb\u898b\u7a4d\u66f8\u30923\u70b9\u30bb\u30c3\u30c8\u3067AI\u304c\u81ea\u52d5\u751f\u6210\u3002\n\u76f8\u624b\u4f01\u696d\u540d\u30fb\u696d\u7a2e\u30fb\u63d0\u6848\u5185\u5bb9\u3092\u5165\u529b\u3059\u308b\u3060\u3051\u3002',sz=16,c=IND,b=True,a=PP_ALIGN.CENTER)
items=[('\u2709\ufe0f \u55b6\u696d\u30e1\u30fc\u30eb','\u696d\u7a2e\u00d7\u63d0\u6848\u5185\u5bb9\u306b\u6700\u9069\u5316\u3055\u308c\u305f\n\u30d1\u30fc\u30bd\u30ca\u30e9\u30a4\u30ba\u30e1\u30fc\u30eb\u3092\u5373\u751f\u6210',IND),
('\U0001f4c4 \u63d0\u6848\u66f8','\u6848\u4ef6\u60c5\u5831\u3092\u5165\u529b\u3059\u308b\u3060\u3051\u3067\n\u30d7\u30ed\u54c1\u8cea\u306e\u63d0\u6848\u66f8\u3092\u81ea\u52d5\u751f\u6210',GR),
('\U0001f4b0 \u898b\u7a4d\u66f8','\u9805\u76ee\u5165\u529b + AI\u88dc\u5b8c\u3067\n\u6b63\u78ba\u306a\u898b\u7a4d\u66f8\u3092\u5373\u5ea7\u306bPDF\u51fa\u529b',WN)]
for i,(t,d,c) in enumerate(items):
    x=M+Inches(i*3.7); y=Inches(2.9)
    rc(s,x,y,Inches(3.3),Inches(1.8),f=LBG,ln=c,lw=Pt(2))
    tx(s,x+Inches(0.2),y+Inches(0.15),Inches(2.9),Inches(0.35),t,sz=16,c=c,b=True)
    tx(s,x+Inches(0.2),y+Inches(0.6),Inches(2.9),Inches(1.0),d,sz=14,c=DN)

# 4. TARGET
s=prs.slides.add_slide(prs.slide_layouts[6]); hd(s,'03','\u30bf\u30fc\u30b2\u30c3\u30c8\u30e6\u30fc\u30b6\u30fc')
personas=[('\U0001f454 \u7530\u4e2d\uff0828\u6b73\uff09','IT\u4f01\u696d\u306e\u6cd5\u4eba\u55b6\u696d','\u30e1\u30fc\u30eb\u4f5c\u6210\u306b\u8ffd\u308f\u308c\n1\u65e520\u793e\u304c\u9650\u754c\u219250\u793e\u3078'),
('\U0001f4bb \u4f50\u85e4\uff0835\u6b73\uff09','\u30d5\u30ea\u30fc\u30e9\u30f3\u30b9\u30a8\u30f3\u30b8\u30cb\u30a2','\u55b6\u696d\u304c\u82e6\u624b\u3002\u63d0\u6848\u66f8\u306b\u534a\u65e5\n\u219230\u5206\u21925\u5206\u306b\u77ed\u7e2e'),
('\U0001f4ca \u9234\u6728\uff0842\u6b73\uff09','\u88fd\u9020\u696d\u306e\u55b6\u696d\u90e8\u9577','\u30c1\u30fc\u30e05\u540d\u306e\u54c1\u8cea\u30d0\u30e9\u30d0\u30e9\n\u2192\u5168\u54e1\u306e\u53d7\u6ce8\u738715%\u5411\u4e0a')]
for i,(icon,role,pain) in enumerate(personas):
    x=M+Inches(i*3.7); y=Inches(1.5)
    rc(s,x,y,Inches(3.3),Inches(2.8),f=LBG,ln=BD); rc(s,x,y,Inches(3.3),Pt(5),f=IND)
    tx(s,x+Inches(0.2),y+Inches(0.2),Inches(2.9),Inches(0.35),icon,sz=16,c=DN,b=True)
    tx(s,x+Inches(0.2),y+Inches(0.6),Inches(2.9),Inches(0.3),role,sz=14,c=IND,b=True)
    tx(s,x+Inches(0.2),y+Inches(1.05),Inches(2.9),Inches(1.2),pain,sz=14,c=SL)

# 5. COMPETITIVE ADVANTAGE
s=prs.slides.add_slide(prs.slide_layouts[6]); hd(s,'04','\u7af6\u4e89\u512a\u4f4d\u6027 \u2015 6\u3064\u306e\u5dee\u5225\u5316')
diffs=[('\u55b6\u696d3\u70b9\u30bb\u30c3\u30c8','\u30e1\u30fc\u30eb\u30fb\u63d0\u6848\u66f8\u30fb\u898b\u7a4d\u66f8\u3092\n1\u30c4\u30fc\u30eb\u3067\u751f\u6210\u3059\u308b\u552f\u4e00\u306e\u30b5\u30fc\u30d3\u30b9'),
('\u696d\u7a2e\u5225AI\u6700\u9069\u5316','\u4f7f\u3046\u307b\u3069\u30c6\u30f3\u30d7\u30ec\u30fc\u30c8\u7cbe\u5ea6\u304c\u5411\u4e0a\n\u30c7\u30fc\u30bf\u30e2\u30fc\u30c8\u304c\u5f62\u6210\u3055\u308c\u308b'),
('\u65e5\u672c\u8a9e\u55b6\u696d\u7279\u5316','\u656c\u8a9e\u30fb\u5b63\u7bc0\u306e\u6328\u62f6\u30fb\u696d\u754c\u7528\u8a9e\n\u65e5\u672c\u306e\u30d3\u30b8\u30cd\u30b9\u6163\u7fd2\u306b\u5b8c\u5168\u5bfe\u5fdc'),
('\u5706\u4fa1\u683c \u00a51,980\u301c','\u65e2\u5b58SFA/MA\u306e1/10\u4ee5\u4e0b\u306e\u4fa1\u683c\n\u500b\u4eba\u30fb\u30d5\u30ea\u30fc\u30e9\u30f3\u30b9\u3067\u3082\u5c0e\u5165\u53ef\u80fd'),
('\u5373\u65e5\u5c0e\u5165\u30fb\u5373\u65e5\u52b9\u679c','\u30a2\u30ab\u30a6\u30f3\u30c8\u4f5c\u6210\u304b\u30895\u5206\u3067\n\u6700\u521d\u306e\u55b6\u696d\u30e1\u30fc\u30eb\u3092\u751f\u6210'),
('\u30c7\u30fc\u30bf\u30e2\u30fc\u30c8','\u696d\u7a2e\u5225\u30c6\u30f3\u30d7\u30ec\u30fc\u30c8\u304c\u84c4\u7a4d\u3055\u308c\n\u5f8c\u767a\u304c\u8ffd\u3044\u3064\u3051\u306a\u3044\u969c\u58c1\u306b')]
for i,(t,d) in enumerate(diffs):
    col=i%3; row=i//3; x=M+Inches(col*3.7); y=Inches(1.5)+Inches(row*2.2)
    rc(s,x,y,Inches(3.3),Inches(1.8),f=LBG,ln=BD); rc(s,x,y,Pt(5),Inches(1.8),f=IND)
    tx(s,x+Inches(0.2),y+Inches(0.12),Inches(2.9),Inches(0.3),'\u2714 '+t,sz=15,c=IND,b=True)
    tx(s,x+Inches(0.2),y+Inches(0.5),Inches(2.9),Inches(1.0),d,sz=14,c=SL)

# 6. PRICING
s=prs.slides.add_slide(prs.slide_layouts[6]); hd(s,'05','\u53ce\u76ca\u30e2\u30c7\u30eb \u2015 3\u30d7\u30e9\u30f3')
plans=[('Starter','\u00a51,980/\u6708','\u500b\u4eba\u30fb\u30d5\u30ea\u30fc\u30e9\u30f3\u30b9','\u30e1\u30fc\u30eb50\u901a/\u6708\n\u63d0\u6848\u66f85\u4ef6/\u6708\n\u898b\u7a4d\u66f85\u4ef6/\u6708',SL),
('Pro','\u00a53,980/\u6708','\u55b6\u696d\u30de\u30f3\u500b\u4eba','\u30e1\u30fc\u30eb200\u901a/\u6708\n\u63d0\u6848\u66f820\u4ef6/\u6708\n\u6210\u679c\u30c8\u30e9\u30c3\u30ad\u30f3\u30b0',IND),
('Team','\u00a54,980/\u4eba/\u6708','\u30c1\u30fc\u30e0\uff083\u540d\u301c\uff09','\u7121\u5236\u9650\u751f\u6210\n\u30c1\u30fc\u30e0\u7ba1\u7406\n\u5206\u6790\u30c0\u30c3\u30b7\u30e5\u30dc\u30fc\u30c9',GR)]
for i,(name,price,target,features,clr) in enumerate(plans):
    x=M+Inches(i*3.7); y=Inches(1.5); cw=Inches(3.3)
    rc(s,x,y,cw,Inches(3.5),f=LBG,ln=clr,lw=Pt(2)); rc(s,x,y,cw,Pt(5),f=clr)
    tx(s,x+Inches(0.2),y+Inches(0.2),cw-Inches(0.4),Inches(0.3),name,sz=18,c=DN,b=True)
    tx(s,x+Inches(0.2),y+Inches(0.55),cw-Inches(0.4),Inches(0.4),price,sz=28,c=clr,b=True)
    tx(s,x+Inches(0.2),y+Inches(1.05),cw-Inches(0.4),Inches(0.25),target,sz=14,c=SL)
    tx(s,x+Inches(0.2),y+Inches(1.5),cw-Inches(0.4),Inches(1.5),features,sz=14,c=DN)
# BEP
rc(s,M,Inches(5.3),CW,Inches(0.6),f=HBG,ln=IND,lw=Pt(2))
tx(s,M+Inches(0.3),Inches(5.35),CW-Inches(0.6),Inches(0.5),'\u640d\u76ca\u5206\u5c90\u70b9: \u305f\u3063\u305f15\u30e6\u30fc\u30b6\u30fc\u3067\u9ed2\u5b57\u5316',sz=16,c=IND,b=True,a=PP_ALIGN.CENTER)

# 7. UNIT ECONOMICS
s=prs.slides.add_slide(prs.slide_layouts[6]); hd(s,'06','\u30e6\u30cb\u30c3\u30c8\u30a8\u30b3\u30ce\u30df\u30af\u30b9')
metrics=[('ARPU','\u00a53,500','\u6708\u984d\u5e73\u5747\u5358\u4fa1'),('LTV','\u00a542,000','\u9867\u5ba2\u751f\u6daf\u4fa1\u5024'),
('CAC','\u00a55,000','\u9867\u5ba2\u7372\u5f97\u30b3\u30b9\u30c8'),('LTV/CAC','8.4\u500d','\u6295\u8cc7\u52b9\u7387')]
for i,(name,val,desc) in enumerate(metrics):
    x=M+Inches(i*2.75); y=Inches(1.5)
    rc(s,x,y,Inches(2.5),Inches(1.8),f=LBG,ln=BD)
    tx(s,x,y+Inches(0.15),Inches(2.5),Inches(0.25),name,sz=14,c=SL,a=PP_ALIGN.CENTER)
    tx(s,x,y+Inches(0.45),Inches(2.5),Inches(0.6),val,sz=32,c=IND,b=True,a=PP_ALIGN.CENTER)
    tx(s,x,y+Inches(1.15),Inches(2.5),Inches(0.3),desc,sz=14,c=SL,a=PP_ALIGN.CENTER)
rc(s,M,Inches(3.7),CW,Inches(0.6),f=RGBColor(209,250,229),ln=GR,lw=Pt(2))
tx(s,M+Inches(0.3),Inches(3.75),CW-Inches(0.6),Inches(0.5),'\u56fa\u5b9a\u8cbb\u6708\u984d\u00a545,000\u306e\u307f\u3002\u5c11\u6570\u30e6\u30fc\u30b6\u30fc\u3067\u5373\u5ea7\u306b\u9ed2\u5b57\u8d85\u3048',sz=15,c=GR,b=True,a=PP_ALIGN.CENTER)

# 8. MARKET SIZE
s=prs.slides.add_slide(prs.slide_layouts[6]); hd(s,'07','\u5e02\u5834\u898f\u6a21')
mkt=[('TAM','\u7d043,000\u5104\u5186','\u65e5\u672c\u56fd\u5185\u306e\u55b6\u696d\u652f\u63f4\u30c4\u30fc\u30eb\u5e02\u5834\u5168\u4f53'),
('SAM','\u7d04300\u5104\u5186','\u4e2d\u5c0f\u4f01\u696d\u30fb\u500b\u4eba\u5411\u3051\n\u55b6\u696d\u30c9\u30ad\u30e5\u30e1\u30f3\u30c8\u81ea\u52d5\u5316'),
('SOM','\u7d043\u5104\u5186','\u521d\u671f\u306b\u7372\u5f97\u53ef\u80fd\u306a\u5e02\u5834\nIT\u7cfb\u30d5\u30ea\u30fc\u30e9\u30f3\u30b9 + \u5c0f\u898f\u6a21\u55b6\u696d')]
for i,(name,val,desc) in enumerate(mkt):
    x=M+Inches(i*3.7); y=Inches(1.5)
    rc(s,x,y,Inches(3.3),Inches(2.5),f=LBG,ln=IND,lw=Pt(2))
    tx(s,x+Inches(0.2),y+Inches(0.15),Inches(2.9),Inches(0.3),name,sz=16,c=IND,b=True)
    tx(s,x+Inches(0.2),y+Inches(0.55),Inches(2.9),Inches(0.5),val,sz=28,c=DN,b=True)
    tx(s,x+Inches(0.2),y+Inches(1.2),Inches(2.9),Inches(1.0),desc,sz=14,c=SL)

# 9. DATA MOAT
s=prs.slides.add_slide(prs.slide_layouts[6]); hd(s,'08','\u30c7\u30fc\u30bf\u30e2\u30fc\u30c8')
rc(s,M,Inches(1.5),CW,Inches(1.0),f=HBG,ln=IND,lw=Pt(2))
tx(s,M+Inches(0.3),Inches(1.6),CW-Inches(0.6),Inches(0.8),'\u696d\u7a2e\u5225\u55b6\u696d\u30c6\u30f3\u30d7\u30ec\u30fc\u30c8\u306e\u6700\u9069\u5316\u30c7\u30fc\u30bf\u304c\u84c4\u7a4d\u3055\u308c\u308b\u307b\u3069\n\u7cbe\u5ea6\u304c\u5411\u4e0a\u3057\u3001\u5f8c\u767a\u304c\u8ffd\u3044\u3064\u3051\u306a\u3044\u969c\u58c1\u306b\u306a\u308b',sz=16,c=IND,b=True,a=PP_ALIGN.CENTER)
steps=[('\u30e6\u30fc\u30b6\u30fc\u304c\n\u55b6\u696d\u30e1\u30fc\u30eb\u751f\u6210',IND),('\u53d7\u6ce8/\u5931\u6ce8\u306e\n\u7d50\u679c\u30c7\u30fc\u30bf\u84c4\u7a4d',GR),
('\u696d\u7a2e\u5225\u30c6\u30f3\u30d7\u30ec\u304c\n\u81ea\u52d5\u6700\u9069\u5316',WN),('\u65b0\u898f\u30e6\u30fc\u30b6\u30fc\u306e\n\u7cbe\u5ea6\u304c\u521d\u65e5\u304b\u3089\u9ad8\u3044',ER)]
for i,(label,clr) in enumerate(steps):
    x=Inches(1.4)+Inches(i*2.8); y=Inches(3.0)
    rc(s,x,y,Inches(2.3),Inches(1.2),f=LBG,ln=clr,lw=Pt(2))
    tx(s,x,y+Inches(0.2),Inches(2.3),Inches(0.8),label,sz=14,c=DN,b=True,a=PP_ALIGN.CENTER)
    if i<3: tx(s,x+Inches(2.35),y+Inches(0.3),Inches(0.4),Inches(0.4),'\u2192',sz=20,c=clr,b=True,a=PP_ALIGN.CENTER)

# 10. ROADMAP
s=prs.slides.add_slide(prs.slide_layouts[6]); hd(s,'09','\u30ed\u30fc\u30c9\u30de\u30c3\u30d7')
phases=[('2026\u5e744\u6708','MVP','\u81ea\u793e\u55b6\u696d\u3067\u4f7f\u7528\u958b\u59cb\n\u30c8\u30e9\u30af\u30b7\u30e7\u30f3\u30c7\u30fc\u30bf\u53d6\u5f97',IND),
('2026\u5e745-6\u6708','v1.0','\u6709\u6599\u30d7\u30e9\u30f3\u958b\u59cb\n\u696d\u7a2e\u5225\u30c6\u30f3\u30d7\u30ec10\u7a2e / LP\u516c\u958b',GR),
('2026\u5e748\u6708\u301c','v2.0','Team\u30d7\u30e9\u30f3 / CRM\u9023\u643a\nA/B\u30c6\u30b9\u30c8\u81ea\u52d5\u5316',WN),
('2027\u5e74\u301c','v3.0','\u30de\u30fc\u30b1\u30c3\u30c8\u30d7\u30ec\u30a4\u30b9 / API\u63d0\u4f9b\nKarakuri\u7d71\u5408 / \u591a\u8a00\u8a9e',ER)]
for i,(period,ver,desc,clr) in enumerate(phases):
    x=M+Inches(i*2.75); y=Inches(1.5)
    rc(s,x,y,Inches(2.5),Inches(3.0),f=LBG,ln=clr,lw=Pt(2)); rc(s,x,y,Inches(2.5),Pt(5),f=clr)
    tx(s,x+Inches(0.15),y+Inches(0.15),Inches(2.2),Inches(0.25),period,sz=13,c=SL)
    tx(s,x+Inches(0.15),y+Inches(0.4),Inches(2.2),Inches(0.35),ver,sz=22,c=clr,b=True)
    tx(s,x+Inches(0.15),y+Inches(0.9),Inches(2.2),Inches(1.8),desc,sz=14,c=DN)

# 11. KPI
s=prs.slides.add_slide(prs.slide_layouts[6]); hd(s,'10','KPI')
rc(s,M,Inches(1.5),CW,Inches(0.45),f=DN)
cols=[(M,3.0,'KPI'),(Inches(4.2),2.2,'6\u30f6\u6708'),(Inches(6.6),2.2,'12\u30f6\u6708'),(Inches(9.0),2.2,'24\u30f6\u6708')]
for x,w,t in cols: tx(s,x,Inches(1.53),Inches(w),Inches(0.4),t,sz=14,c=WH,b=True,a=PP_ALIGN.CENTER)
rows=[('\u6709\u6599\u30e6\u30fc\u30b6\u30fc','50\u4eba','200\u4eba','1,000\u4eba'),('MRR','\u00a5175,000','\u00a5700,000','\u00a53,500,000'),
('ARR','\u2015','\u00a58,400,000','\u00a542,000,000'),('\u30c1\u30e3\u30fc\u30f3\u30ec\u30fc\u30c8','8%\u4ee5\u4e0b','5%\u4ee5\u4e0b','3%\u4ee5\u4e0b')]
for i,(label,v1,v2,v3) in enumerate(rows):
    y=Inches(1.95)+Inches(i*0.5)
    rc(s,M,y,CW,Inches(0.5),f=LBG if i%2==0 else WH,ln=BD)
    tx(s,M+Inches(0.1),y+Inches(0.07),Inches(2.8),Inches(0.35),label,sz=14,c=DN,b=True)
    for j,v in enumerate([v1,v2,v3]):
        tx(s,Inches(4.2)+Inches(j*2.4),y+Inches(0.07),Inches(2.2),Inches(0.35),v,sz=14,c=DN,b=True,a=PP_ALIGN.CENTER)

# 12. MVP SUMMARY
s=prs.slides.add_slide(prs.slide_layouts[6]); hd(s,'11','MVP\u4ed5\u69d8\u30b5\u30de\u30ea\u30fc')
specs=[('Must Have\u6a5f\u80fd','8\u6a5f\u80fd','\u8a8d\u8a3c / \u30e1\u30fc\u30eb\u751f\u6210 / \u63d0\u6848\u66f8\n\u898b\u7a4d\u66f8 / \u5c65\u6b74 / PDF / \u30c6\u30f3\u30d7\u30ec / \u30b3\u30d4\u30fc'),
('\u753b\u9762\u6570','8\u753b\u9762','\u30ed\u30b0\u30a4\u30f3 / \u30c0\u30c3\u30b7\u30e5\u30dc\u30fc\u30c9\n\u30e1\u30fc\u30eb / \u63d0\u6848\u66f8 / \u898b\u7a4d\u66f8\n\u8a73\u7d30 / \u8a2d\u5b9a'),
('DB\u30c6\u30fc\u30d6\u30eb','3\u30c6\u30fc\u30d6\u30eb','profiles / documents\nindustry_templates'),
('\u958b\u767a\u671f\u9593','3\u9031\u9593','W1: \u57fa\u76e4+\u30e1\u30fc\u30eb\u751f\u6210\nW2: \u63d0\u6848\u66f8+\u898b\u7a4d\u66f8+\u5c65\u6b74\nW3: PDF+\u30c6\u30f3\u30d7\u30ec+\u30ea\u30ea\u30fc\u30b9')]
for i,(name,val,desc) in enumerate(specs):
    col=i%2; row=i//2; x=M+Inches(col*5.6); y=Inches(1.5)+Inches(row*2.4)
    rc(s,x,y,Inches(5.0),Inches(2.0),f=LBG,ln=BD); rc(s,x,y,Pt(5),Inches(2.0),f=IND)
    tx(s,x+Inches(0.2),y+Inches(0.1),Inches(2.5),Inches(0.3),name,sz=14,c=SL)
    tx(s,x+Inches(0.2),y+Inches(0.4),Inches(2.5),Inches(0.4),val,sz=24,c=IND,b=True)
    tx(s,x+Inches(0.2),y+Inches(0.9),Inches(4.5),Inches(1.0),desc,sz=14,c=DN)

out=os.path.join(os.path.dirname(os.path.abspath(__file__)),'SlideShot_\u8a2d\u8a08\u66f8.pptx')
prs.save(out); print(f'Saved: {out}')
